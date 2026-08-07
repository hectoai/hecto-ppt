"""헥토 슬라이드 자료가 DESIGN.md 스펙에 맞는지 검사한다.

사용법:
    python scripts/verify_pptx.py <deck.pptx>

스펙 위반이 하나라도 있으면 exit 1.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

TOL = 0.005  # inch. EMU 반올림 오차 흡수용

CANVAS_W, CANVAS_H = 13.333, 7.5
FONT = "Pretendard"

# DESIGN.md > Colors
PALETTE = {
    "FF6013", "FFE7DC", "FFBB93",          # brand
    "FFFFFF", "D0CECF", "E8E9EC", "EFF0F2",  # surface
    "000000", "595959", "838383", "A6AAA9",  # text
    "15B886", "FF0000", "00A1FF",          # semantic
}

# 셰브론 헤더 앵커 (DESIGN.md:785). z-order 순서 그대로.
CHEVRON = [
    ("context-bar",    "rect",  2.95, 0.25, 10.03, 0.59),
    ("chevron-white",  "delay", 2.95, 0.25, 0.34, 0.59),
    ("chapter-pill",   "rect",  0.40, 0.25, 2.30, 0.59),
    ("chevron-orange", "delay", 2.68, 0.25, 0.43, 0.59),
]
WORDMARK = (11.80, 0.353, 1.06, 0.38)

SOURCE_Y, SOURCE_W = 7.10, 11.50  # DESIGN.md:125 Vertical rhythm 표
PAGENUM_X, PAGENUM_Y = 12.13, 7.10
END_MARK_X, END_MARK_Y = 0.40, 3.00
TABLE_ROW_H = 0.42          # DESIGN.md:320 data-table
BRAND_ORANGE = "FF6013"
ASPECT_TOL = 0.005          # 워드마크 허용 왜곡 0.5%. 픽셀 반올림 잔차는 통과시킨다

# 고정 보일러플레이트 (DESIGN.md:460). 줄바꿈은 공백으로 정규화해 비교한다.
CONFIDENTIAL = ("Confidential and proprietary Any use of this material without "
                "specific permission of Hecto is strictly prohibited")


class Report:
    def __init__(self):
        self.failures = []

    def fail(self, rule, detail):
        self.failures.append((rule, detail))

    def emit(self, rules):
        by_rule = {}
        for rule, detail in self.failures:
            by_rule.setdefault(rule, []).append(detail)
        for rule in rules:
            if rule in by_rule:
                print(f"  FAIL  {rule}")
                for d in by_rule[rule]:
                    print(f"          {d}")
            else:
                print(f"  ok    {rule}")
        return len(self.failures)


def inches(v):
    return Emu(v).inches if v is not None else None


def box(sh):
    return (inches(sh.left), inches(sh.top), inches(sh.width), inches(sh.height))


def near(a, b):
    return a is not None and abs(a - b) <= TOL


def is_delay(sh):
    try:
        return "FLOWCHART_DELAY" in str(sh.auto_shape_type)
    except (ValueError, AttributeError):
        return False


def is_picture(sh):
    return "PICTURE" in str(sh.shape_type)


def solid_fill(sh):
    try:
        if sh.fill.type == 1:
            return str(sh.fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def runs(sh):
    if not sh.has_text_frame:
        return
    for para in sh.text_frame.paragraphs:
        for r in para.runs:
            yield r


def text_of(sh):
    return sh.text_frame.text.strip() if sh.has_text_frame else ""


def is_micro(sh):
    """typography.micro(10pt) 급인지. 번호 마커(40pt)와 페이지 번호를 가른다."""
    sizes = [r.font.size.pt for r in runs(sh) if r.font.size]
    return bool(sizes) and max(sizes) <= 12.0


def check_canvas(prs, rep):
    w, h = inches(prs.slide_width), inches(prs.slide_height)
    if not (near(w, CANVAS_W) and near(h, CANVAS_H)):
        rep.fail("R1 캔버스 13.333x7.5",
                 f"실측 {w:.3f} x {h:.3f}")


def check_chevron(idx, slide, rep):
    """콘텐츠 슬라이드의 셰브론 헤더 좌표와 z-order."""
    got = []
    for sh in slide.shapes:
        if is_picture(sh):
            continue
        fill = solid_fill(sh)
        if fill is None:
            continue
        got.append((sh, box(sh), is_delay(sh)))
        if len(got) == len(CHEVRON):
            break

    if len(got) < len(CHEVRON):
        rep.fail("R2 셰브론 헤더", f"slide {idx}: 헤더 도형 {len(got)}/{len(CHEVRON)}개만 발견")
        return

    for (name, kind, x, y, w, h), (sh, (gx, gy, gw, gh), delay) in zip(CHEVRON, got):
        if not all([near(gx, x), near(gy, y), near(gw, w), near(gh, h)]):
            rep.fail("R2 셰브론 헤더",
                     f"slide {idx} {name}: 스펙 ({x} {y} {w} {h}) / 실측 "
                     f"({gx:.3f} {gy:.3f} {gw:.3f} {gh:.3f})")
        if (kind == "delay") != delay:
            rep.fail("R2 셰브론 헤더",
                     f"slide {idx} {name}: 도형 종류 불일치 (z-order 어긋남 가능)")

    pics = [sh for sh in slide.shapes if is_picture(sh)]
    if not pics:
        rep.fail("R2 셰브론 헤더", f"slide {idx}: 워드마크 없음")
    else:
        gx, gy, gw, gh = box(pics[0])
        x, y, w, h = WORDMARK
        if not all([near(gx, x), near(gy, y), near(gw, w), near(gh, h)]):
            rep.fail("R2 셰브론 헤더",
                     f"slide {idx} wordmark: 스펙 ({x} {y} {w} {h}) / 실측 "
                     f"({gx:.3f} {gy:.3f} {gw:.3f} {gh:.3f})")


def check_palette_and_font(idx, slide, rep):
    for sh in slide.shapes:
        fill = solid_fill(sh)
        if fill and fill not in PALETTE:
            rep.fail("R3 팔레트", f"slide {idx}: 미등록 채움색 #{fill}")
        for r in runs(sh):
            if r.font.name and r.font.name != FONT:
                rep.fail("R4 폰트", f"slide {idx}: '{r.font.name}' (Pretendard 아님)")
            try:
                c = str(r.font.color.rgb).upper()
            except Exception:
                continue
            if c not in PALETTE:
                rep.fail("R3 팔레트", f"slide {idx}: 미등록 글자색 #{c}")


def check_footer(idx, slide, rep):
    for sh in slide.shapes:
        t = text_of(sh)
        if t.startswith("출처:"):
            y, w = inches(sh.top), inches(sh.width)
            if not near(y, SOURCE_Y):
                rep.fail("R5 출처 캡션 위치",
                         f"slide {idx}: y 실측 {y:.3f} ({y - SOURCE_Y:+.3f})")
            if not near(w, SOURCE_W):
                rep.fail("R5 출처 캡션 위치",
                         f"slide {idx}: w 실측 {w:.3f} (페이지 번호와 겹칠 수 있음)")
            if "자체 분석" in t:
                rep.fail("R5 출처 캡션 위치",
                         f"slide {idx}: '자체 분석'은 출처로 쓰지 않는다")
        elif t.isdigit() and is_micro(sh):
            x, y = inches(sh.left), inches(sh.top)
            if not (near(x, PAGENUM_X) and near(y, PAGENUM_Y)):
                rep.fail("R6 페이지 번호 위치",
                         f"slide {idx}: 실측 ({x:.3f}, {y:.3f})")


def check_cover(slide, rep):
    """컨피덴셜 블록은 verbatim (DESIGN.md:460). 번역·축약·마침표 추가 전부 금지."""
    for sh in slide.shapes:
        t = text_of(sh)
        if not t.startswith("Confidential"):
            continue
        if " ".join(t.split()) != CONFIDENTIAL:
            rep.fail("R9 커버 보일러플레이트", f"원문과 다름: {t!r}")
        return
    rep.fail("R9 커버 보일러플레이트", "커버에 컨피덴셜 블록이 없음")


def check_tables(idx, slide, rep):
    """데이터 표 (DESIGN.md:317).

    행 높이 0.42". 강조 열 헤더는 **글자색**이 brand-orange일 뿐 배경은 다른 헤더와
    같은 wash다. 배경까지 오렌지로 채우면 표 절반이 시그니처 오렌지가 되어
    페이지의 오렌지 예산을 통째로 삼킨다.
    """
    for sh in slide.shapes:
        if not getattr(sh, "has_table", False):
            continue
        for r, row in enumerate(sh.table.rows):
            h = inches(row.height)
            if not near(h, TABLE_ROW_H):
                rep.fail("R10 데이터 표",
                         f"slide {idx}: {r + 1}행 높이 {h:.3f} (스펙 {TABLE_ROW_H})")
        for c, cell in enumerate(sh.table.rows[0].cells):
            try:
                fill = str(cell.fill.fore_color.rgb).upper()
            except Exception:
                continue
            if fill == BRAND_ORANGE:
                rep.fail("R10 데이터 표",
                         f"slide {idx}: 헤더 {c + 1}열 배경이 brand-orange "
                         f"(강조는 글자색으로만 준다)")


def check_wordmark_aspect(idx, slide, rep):
    """워드마크가 눌리지 않았는가 (DESIGN.md Asset Processing).

    스펙은 "Never modify the wordmark's proportions"라고 못 박는다. 그 조항은
    컨테이너 종횡비를 embed된 PNG의 실제 종횡비와 대조해야만 강제할 수 있다.
    에셋을 갈아끼워도 여기서 걸린다.
    """
    for sh in slide.shapes:
        if not is_picture(sh):
            continue
        w, h = inches(sh.width), inches(sh.height)
        if not w or not h:
            continue
        try:
            px_w, px_h = sh.image.size
        except Exception:
            continue
        drift = (w / h) / (px_w / px_h) - 1
        if abs(drift) > ASPECT_TOL:
            rep.fail("R11 워드마크 비율",
                     f"slide {idx}: 컨테이너 {w:.3f}x{h:.3f} ({w / h:.4f}) vs "
                     f"이미지 {px_w}x{px_h} ({px_w / px_h:.4f}) → {drift * 100:+.2f}% 왜곡")


def check_bounds(idx, slide, rep):
    for sh in slide.shapes:
        x, y, w, h = box(sh)
        if None in (x, y, w, h):
            continue
        if x < -TOL or y < -TOL or x + w > CANVAS_W + TOL or y + h > CANVAS_H + TOL:
            rep.fail("R7 캔버스 이탈",
                     f"slide {idx}: ({x:.3f} {y:.3f} {w:.3f} {h:.3f})")


def check_end_slide(slides, rep):
    """모든 자료는 slide-end로 닫는다 (DESIGN.md:792)."""
    last = slides[-1]
    mark = None
    for sh in last.shapes:
        if text_of(sh) in ("E. O. D.", "끝."):
            mark = sh
            break
    if mark is None:
        rep.fail("R8 E.O.D. 종료 슬라이드",
                 f"마지막 슬라이드({len(slides)})에 종료 마크가 없음")
        return

    x, y = inches(mark.left), inches(mark.top)
    if not (near(x, END_MARK_X) and near(y, END_MARK_Y)):
        rep.fail("R8 E.O.D. 종료 슬라이드",
                 f"종료 마크 위치: 스펙 ({END_MARK_X}, {END_MARK_Y}) / 실측 ({x:.3f}, {y:.3f})")

    expected_pt = 72.0 if text_of(mark) == "E. O. D." else 88.0
    for r in runs(mark):
        if r.font.size and abs(r.font.size.pt - expected_pt) > 0.1:
            rep.fail("R8 E.O.D. 종료 슬라이드",
                     f"종료 마크 크기: 스펙 {expected_pt}pt / 실측 {r.font.size.pt}pt")


RULES = [
    "R1 캔버스 13.333x7.5",
    "R2 셰브론 헤더",
    "R3 팔레트",
    "R4 폰트",
    "R5 출처 캡션 위치",
    "R6 페이지 번호 위치",
    "R7 캔버스 이탈",
    "R8 E.O.D. 종료 슬라이드",
    "R9 커버 보일러플레이트",
    "R10 데이터 표",
    "R11 워드마크 비율",
]


def main(path):
    prs = Presentation(path)
    slides = list(prs.slides)
    rep = Report()

    check_canvas(prs, rep)
    check_cover(slides[0], rep)
    for i, slide in enumerate(slides, start=1):
        if any(is_delay(sh) for sh in slide.shapes):
            check_chevron(i, slide, rep)
        check_footer(i, slide, rep)  # 섹션 디바이더도 페이지 번호를 갖는다
        check_palette_and_font(i, slide, rep)
        check_tables(i, slide, rep)
        check_wordmark_aspect(i, slide, rep)
        check_bounds(i, slide, rep)
    check_end_slide(slides, rep)

    print(f"검사 대상: {path}")
    print(f"슬라이드 {len(slides)}장\n")
    n = rep.emit(RULES)
    print()
    if n:
        print(f"위반 {n}건")
        return 1
    print("전부 통과")
    return 0


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(__doc__)
        sys.exit(2)
    sys.exit(main(sys.argv[1]))
